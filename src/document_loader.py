import os
import json
from pathlib import Path
from typing import List, Dict, Any
from datetime import datetime

from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    CSVLoader,
    UnstructuredMarkdownLoader,
    BSHTMLLoader
)

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class MultiFormatDocumentLoader:
    """
    Cargador de documentos multi-formato corporativo.
    Soporta PDF, Word (DOCX), Excel (XLSX/CSV), PowerPoint (PPTX), Markdown (MD), JSON e HTML.
    Limpia el texto e incluye metadatos enriquecidos (categoría, archivo de origen, tipo, ubicación).
    """

    def __init__(self, categories_map: Dict[str, Any] = None):
        self.categories_map = categories_map or {}

    def _infer_category(self, file_path: Path) -> str:
        """Determina la categoría corporativa basándose en la subcarpeta o nombre del archivo."""
        parent_folder = file_path.parent.name.lower()
        file_name = file_path.name.lower()

        for cat_name, info in self.categories_map.items():
            dir_keyword = info.get("dir", "").lower()
            if dir_keyword and (dir_keyword in parent_folder or dir_keyword in file_name):
                return cat_name
            if cat_name.lower() in parent_folder or cat_name.lower() in file_name:
                return cat_name
                
        return "General / Otro"

    def load_document(self, file_path: str | Path, category_override: str = None) -> List[Document]:
        """Carga un documento individual según su extensión y retorna lista de LangChain Documents."""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"El archivo no existe: {file_path}")

        ext = file_path.suffix.lower()
        category = category_override or self._infer_category(file_path)
        base_metadata = {
            "source": file_path.name,
            "file_path": str(file_path),
            "category": category,
            "file_type": ext.replace(".", ""),
            "ingested_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }

        docs = []

        try:
            if ext == ".pdf":
                docs = self._load_pdf(file_path, base_metadata)
            elif ext in [".docx", ".doc"]:
                docs = self._load_docx(file_path, base_metadata)
            elif ext == ".csv":
                docs = self._load_csv(file_path, base_metadata)
            elif ext in [".xlsx", ".xls"]:
                docs = self._load_excel(file_path, base_metadata)
            elif ext == ".pptx":
                docs = self._load_pptx(file_path, base_metadata)
            elif ext in [".md", ".markdown"]:
                docs = self._load_markdown(file_path, base_metadata)
            elif ext == ".json":
                docs = self._load_json(file_path, base_metadata)
            elif ext in [".html", ".htm"]:
                docs = self._load_html(file_path, base_metadata)
            elif ext == ".txt":
                docs = self._load_text(file_path, base_metadata)
            else:
                print(f"[WARN] Formato no soportado directamente: {ext}. Intentando carga como texto plano.")
                docs = self._load_text(file_path, base_metadata)

        except Exception as e:
            print(f"[ERROR] Error cargando archivo {file_path.name}: {str(e)}")

        return docs

    def _load_pdf(self, path: Path, metadata: dict) -> List[Document]:
        loader = PyPDFLoader(str(path))
        raw_docs = loader.load()
        for d in raw_docs:
            d.metadata.update(metadata)
            d.metadata["page"] = d.metadata.get("page", 0) + 1
        return raw_docs

    def _load_docx(self, path: Path, metadata: dict) -> List[Document]:
        loader = Docx2txtLoader(str(path))
        raw_docs = loader.load()
        for d in raw_docs:
            d.metadata.update(metadata)
        return raw_docs

    def _load_csv(self, path: Path, metadata: dict) -> List[Document]:
        if pd is not None:
            df = pd.read_csv(path).fillna("")
            docs = []
            for idx, row in df.iterrows():
                row_str = " | ".join([f"{col}: {val}" for col, val in row.items()])
                meta = metadata.copy()
                meta["row"] = idx + 1
                docs.append(Document(page_content=row_str, metadata=meta))
            return docs
        else:
            loader = CSVLoader(str(path))
            raw_docs = loader.load()
            for d in raw_docs:
                d.metadata.update(metadata)
            return raw_docs

    def _load_excel(self, path: Path, metadata: dict) -> List[Document]:
        if pd is None:
            raise ImportError("pandas es requerido para procesar archivos Excel (.xlsx/.xls)")
        
        excel_file = pd.ExcelFile(path)
        docs = []
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name).fillna("")
            for idx, row in df.iterrows():
                row_str = f"Hoja '{sheet_name}' - Fila {idx + 1}: " + ", ".join([f"{col}: {val}" for col, val in row.items() if str(val).strip() != ""])
                meta = metadata.copy()
                meta["sheet"] = sheet_name
                meta["row"] = idx + 1
                docs.append(Document(page_content=row_str, metadata=meta))
        return docs

    def _load_pptx(self, path: Path, metadata: dict) -> List[Document]:
        if Presentation is None:
            # Fallback simple si python-pptx no estuviera disponible
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return [Document(page_content=f.read(), metadata=metadata)]
        
        prs = Presentation(path)
        docs = []
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # Incluir notas de la diapositiva si existen
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text.append(f"[Notas del Orador: {notes}]")

            content = "\n".join(slide_text)
            if content.strip():
                meta = metadata.copy()
                meta["slide"] = idx + 1
                docs.append(Document(page_content=f"Diapositiva {idx + 1}:\n" + content, metadata=meta))
        return docs

    def _load_markdown(self, path: Path, metadata: dict) -> List[Document]:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return [Document(page_content=text, metadata=metadata)]

    def _load_json(self, path: Path, metadata: dict) -> List[Document]:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            data = json.load(f)

        docs = []
        if isinstance(data, list):
            for idx, item in enumerate(data):
                content = json.dumps(item, ensure_ascii=False, indent=2)
                meta = metadata.copy()
                meta["item_index"] = idx + 1
                docs.append(Document(page_content=content, metadata=meta))
        elif isinstance(data, dict):
            for key, val in data.items():
                content = f"Sección '{key}':\n" + json.dumps(val, ensure_ascii=False, indent=2)
                meta = metadata.copy()
                meta["key"] = key
                docs.append(Document(page_content=content, metadata=meta))
        else:
            docs.append(Document(page_content=str(data), metadata=metadata))
        return docs

    def _load_html(self, path: Path, metadata: dict) -> List[Document]:
        loader = BSHTMLLoader(str(path), open_encoding="utf-8")
        raw_docs = loader.load()
        for d in raw_docs:
            d.metadata.update(metadata)
        return raw_docs

    def _load_text(self, path: Path, metadata: dict) -> List[Document]:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return [Document(page_content=text, metadata=metadata)]

    def load_directory(self, dir_path: str | Path) -> List[Document]:
        """Recorre recursivamente un directorio e ingesta todos los documentos compatibles."""
        dir_path = Path(dir_path)
        all_docs = []

        if not dir_path.exists():
            return all_docs

        for file_path in dir_path.rglob("*"):
            if file_path.is_file() and not file_path.name.startswith("."):
                docs = self.load_document(file_path)
                all_docs.extend(docs)

        return all_docs
